"""
LocalRAG AI Knowledge Assistant - Standalone Version

A single-file version with all core functionality for reliable Streamlit Cloud deployment.
"""

import os
import sys
import tempfile
import json
import pickle
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from abc import ABC, abstractmethod
import streamlit as st

# Third-party imports with error handling
missing_packages = []

try:
    import numpy as np
except ImportError as e:
    missing_packages.append("numpy")

try:
    import faiss
except ImportError as e:
    missing_packages.append("faiss-cpu")

try:
    from sentence_transformers import SentenceTransformer
except ImportError as e:
    missing_packages.append("sentence-transformers")

try:
    from groq import Groq
except ImportError as e:
    missing_packages.append("groq")

try:
    import PyPDF2
except ImportError as e:
    missing_packages.append("PyPDF2")

try:
    import pandas as pd
except ImportError as e:
    missing_packages.append("pandas")

try:
    from docx import Document
except ImportError as e:
    missing_packages.append("python-docx")

try:
    from pptx import Presentation
except ImportError as e:
    missing_packages.append("python-pptx")

try:
    import markdown
except ImportError as e:
    missing_packages.append("markdown")

# Load dotenv if available (optional)
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# If any required packages are missing, show error and stop
if missing_packages:
    st.error("❌ Missing required dependencies:")
    for pkg in missing_packages:
        st.error(f"  • {pkg}")
    
    st.info(f"""
    ### 🛠️ To fix this issue:
    
    1. **If running locally**:
       ```bash
       pip install {" ".join(missing_packages)}
       ```
    
    2. **If deploying to Streamlit Cloud**:
       Make sure your `requirements.txt` includes:
       ```
       {"\
       ".join(missing_packages)}
       ```
    """)
    
    st.stop()

# ==================== CORE MODULES ====================

# Document Processor
class DocumentProcessor:
    """Enterprise-grade document processing engine."""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.txt', '.doc', '.docx', '.csv', '.xlsx', '.xls', '.md', '.pptx'}
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
    def process_file(self, file_path: str) -> Dict[str, Any]:
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Extract text based on file type
        text = self._extract_text_by_type(file_path, file_extension)
        
        # Create intelligent chunks
        chunks = self._create_chunks(text)
        
        return {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_type': file_extension,
            'text': text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'file_size': file_path.stat().st_size
        }
    
    def _extract_text_by_type(self, file_path: Path, file_extension: str) -> str:
        """Extract text based on file type."""
        extractors = {
            '.pdf': self._extract_pdf_text,
            '.txt': self._extract_txt_text,
            '.doc': self._extract_docx_text,
            '.docx': self._extract_docx_text,
            '.csv': self._extract_excel_text,
            '.xlsx': self._extract_excel_text,
            '.xls': self._extract_excel_text,
            '.md': self._extract_markdown_text,
            '.pptx': self._extract_pptx_text
        }
        
        extractor = extractors.get(file_extension)
        if not extractor:
            raise ValueError(f"No extractor available for {file_extension}")
            
        return extractor(file_path)
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF {file_path}: {str(e)}")
        return text.strip()
    
    def _extract_txt_text(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        raise Exception(f"Could not decode text file {file_path}")
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading DOCX {file_path}: {str(e)}")
    
    def _extract_excel_text(self, file_path: Path) -> str:
        """Extract text from Excel/CSV file."""
        try:
            if file_path.suffix.lower() == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            text = f"Document: {file_path.name}\n\n"
            text += df.to_string(index=False)
            return text
        except Exception as e:
            raise Exception(f"Error reading Excel/CSV {file_path}: {str(e)}")
    
    def _extract_markdown_text(self, file_path: Path) -> str:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to plain text
            html = markdown.markdown(md_content)
            text = re.sub('<[^<]+?>', '', html)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading Markdown {file_path}: {str(e)}")
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file_path)
            text = f"Presentation: {file_path.name}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"Slide {i}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading PPTX {file_path}: {str(e)}")
    
    def _create_chunks(self, text: str) -> List[Dict[str, Any]]:
        """Create intelligent text chunks with overlap."""
        if not text.strip():
            return []
        
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        chunks = []
        current_chunk = ""
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            if current_size + sentence_size > self.chunk_size and current_chunk:
                chunks.append({
                    'text': current_chunk.strip(),
                    'size': current_size,
                    'chunk_id': len(chunks)
                })
                
                overlap_text = self._get_overlap_text(current_chunk)
                current_chunk = overlap_text + " " + sentence
                current_size = len(current_chunk)
            else:
                current_chunk += " " + sentence if current_chunk else sentence
                current_size += sentence_size + (1 if current_chunk else 0)
        
        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'size': current_size,
                'chunk_id': len(chunks)
            })
        
        return chunks
    
    def _get_overlap_text(self, text: str) -> str:
        """Get intelligent overlap text from chunk end."""
        if len(text) <= self.chunk_overlap:
            return text
        
        overlap_start = len(text) - self.chunk_overlap
        overlap_text = text[overlap_start:]
        
        # Find sentence boundary
        sentence_break = overlap_text.find('. ')
        if sentence_break != -1:
            return overlap_text[sentence_break + 2:]
        
        # Find word boundary
        word_break = overlap_text.find(' ')
        if word_break != -1:
            return overlap_text[word_break + 1:]
        
        return overlap_text

# Vector Store
class VectorStore(ABC):
    """Abstract base class for vector stores."""
    
    @abstractmethod
    def add_documents(self, documents: List[Dict[str, Any]]) -> None:
        pass
    
    @abstractmethod
    def search(self, query: str, k: int = 5, score_threshold: float = 0.0) -> List[Dict[str, Any]]:
        pass
    
    @abstractmethod
    def get_stats(self) -> Dict[str, Any]:
        pass

class FAISSVectorStore(VectorStore):
    """Professional FAISS-based vector store with enterprise features."""
    
    EMBEDDING_MODELS = {
        'all-MiniLM-L6-v2': 384,
        'all-mpnet-base-v2': 768,
        'all-distilroberta-v1': 768,
        'paraphrase-MiniLM-L6-v2': 384,
        'paraphrase-mpnet-base-v2': 768
    }
    
    def __init__(self, 
                 store_path: str = "data/vector_store",
                 embedding_model: str = "all-MiniLM-L6-v2"):
        self.store_path = Path(store_path)
        self.store_path.mkdir(parents=True, exist_ok=True)
        
        self.embedding_model_name = embedding_model
        self.dimension = self.EMBEDDING_MODELS.get(embedding_model, 384)
        
        # Initialize sentence transformer
        st.sidebar.info(f"🔄 Loading embedding model: {embedding_model}")
        self.embedding_model = SentenceTransformer(embedding_model)
        
        # Initialize FAISS index with inner product for cosine similarity
        self.index = faiss.IndexFlatIP(self.dimension)
        
        # Metadata storage
        self.metadata = []
        self.document_map = {}
        
        # Load existing data
        self._load_store()
        st.sidebar.success(f"✅ Vector store initialized with {self.index.ntotal} vectors")
    
    def add_documents(self, documents: List[Dict[str, Any]]) -> None:
        all_texts = []
        all_metadata = []
        
        for doc in documents:
            file_path = doc['file_path']
            file_name = doc['file_name']
            file_type = doc['file_type']
            
            for chunk in doc['chunks']:
                all_texts.append(chunk['text'])
                
                metadata = {
                    'file_path': file_path,
                    'file_name': file_name,
                    'file_type': file_type,
                    'chunk_id': chunk['chunk_id'],
                    'chunk_size': chunk['size'],
                    'text': chunk['text']
                }
                all_metadata.append(metadata)
        
        if not all_texts:
            st.warning("⚠️ No texts to add to vector store")
            return
        
        st.info(f"🔄 Generating embeddings for {len(all_texts)} text chunks...")
        
        # Generate embeddings in batches
        embeddings = self._generate_embeddings_batch(all_texts)
        
        # Add to FAISS index
        self.index.add(embeddings)
        
        # Add metadata
        self.metadata.extend(all_metadata)
        
        # Update document map
        for doc in documents:
            doc_id = doc['file_path']
            self.document_map[doc_id] = {
                'file_name': doc['file_name'],
                'file_type': doc['file_type'],
                'num_chunks': doc['num_chunks'],
                'file_size': doc.get('file_size', 0)
            }
        
        st.success(f"✅ Added {len(all_texts)} chunks to vector store")
        st.sidebar.info(f"📊 Total vectors in store: {self.index.ntotal}")
        
        # Persist changes
        self._save_store()
    
    def search(self, query: str, k: int = 5, score_threshold: float = 0.0) -> List[Dict[str, Any]]:
        if self.index.ntotal == 0:
            return []
        
        # Generate query embedding
        query_embedding = self._generate_embeddings_batch([query])
        
        # Search FAISS index
        scores, indices = self.index.search(query_embedding, min(k, self.index.ntotal))
        
        results = []
        for score, idx in zip(scores[0], indices[0]):
            if idx == -1 or score < score_threshold:
                continue
            
            metadata = self.metadata[idx].copy()
            metadata['similarity_score'] = float(score)
            results.append(metadata)
        
        return results
    
    def get_stats(self) -> Dict[str, Any]:
        return {
            'total_vectors': self.index.ntotal,
            'total_documents': len(self.document_map),
            'embedding_model': self.embedding_model_name,
            'dimension': self.dimension,
            'store_path': str(self.store_path),
            'documents': list(self.document_map.keys()),
            'storage_size_mb': self._get_storage_size()
        }
    
    def delete_document(self, file_path: str) -> bool:
        if file_path not in self.document_map:
            return False
        
        # Find indices to keep
        indices_to_keep = []
        new_metadata = []
        
        for i, metadata in enumerate(self.metadata):
            if metadata['file_path'] != file_path:
                indices_to_keep.append(i)
                new_metadata.append(metadata)
        
        if len(indices_to_keep) == len(self.metadata):
            return False
        
        # Rebuild index
        if indices_to_keep:
            remaining_texts = [meta['text'] for meta in new_metadata]
            embeddings = self._generate_embeddings_batch(remaining_texts)
            
            self.index = faiss.IndexFlatIP(self.dimension)
            self.index.add(embeddings)
        else:
            self.index = faiss.IndexFlatIP(self.dimension)
        
        # Update metadata
        self.metadata = new_metadata
        del self.document_map[file_path]
        
        self._save_store()
        st.sidebar.info(f"🗑️ Deleted document: {file_path}")
        return True
    
    def clear_store(self) -> None:
        self.index = faiss.IndexFlatIP(self.dimension)
        self.metadata = []
        self.document_map = {}
        self._save_store()
        st.sidebar.info("🧹 Vector store cleared")
    
    def _generate_embeddings_batch(self, texts: List[str], batch_size: int = 32) -> np.ndarray:
        """Generate embeddings in batches to manage memory."""
        embeddings = []
        
        for i in range(0, len(texts), batch_size):
            batch_texts = texts[i:i + batch_size]
            batch_embeddings = self.embedding_model.encode(
                batch_texts,
                convert_to_numpy=True,
                normalize_embeddings=True
            )
            embeddings.append(batch_embeddings)
        
        return np.vstack(embeddings)
    
    def _get_storage_size(self) -> float:
        """Get storage size in MB."""
        total_size = 0
        for file_path in self.store_path.glob('*'):
            if file_path.is_file():
                total_size += file_path.stat().st_size
        return round(total_size / (1024 * 1024), 2)
    
    def _save_store(self) -> None:
        """Persist vector store to disk."""
        try:
            # Save FAISS index
            index_path = self.store_path / "faiss_index.bin"
            faiss.write_index(self.index, str(index_path))
            
            # Save metadata
            metadata_path = self.store_path / "metadata.pkl"
            with open(metadata_path, 'wb') as f:
                pickle.dump(self.metadata, f)
            
            # Save document map
            doc_map_path = self.store_path / "document_map.json"
            with open(doc_map_path, 'w') as f:
                json.dump(self.document_map, f, indent=2)
            
            # Save configuration
            config_path = self.store_path / "config.json"
            config = {
                'embedding_model': self.embedding_model_name,
                'dimension': self.dimension,
                'total_vectors': self.index.ntotal,
                'version': '2.0.0'
            }
            with open(config_path, 'w') as f:
                json.dump(config, f, indent=2)
                
        except Exception as e:
            st.error(f"❌ Error saving vector store: {str(e)}")
    
    def _load_store(self) -> None:
        """Load vector store from disk."""
        try:
            index_path = self.store_path / "faiss_index.bin"
            metadata_path = self.store_path / "metadata.pkl"
            doc_map_path = self.store_path / "document_map.json"
            config_path = self.store_path / "config.json"
            
            if not all(p.exists() for p in [index_path, metadata_path, doc_map_path]):
                st.sidebar.info("📁 No existing vector store found, starting fresh")
                return
            
            # Load configuration
            if config_path.exists():
                with open(config_path, 'r') as f:
                    config = json.load(f)
                    
                if config['embedding_model'] != self.embedding_model_name:
                    st.sidebar.warning(f"⚠️ Model mismatch: stored({config['embedding_model']}) vs current({self.embedding_model_name})")
            
            # Load FAISS index
            self.index = faiss.read_index(str(index_path))
            
            # Load metadata
            with open(metadata_path, 'rb') as f:
                self.metadata = pickle.load(f)
            
            # Load document map
            with open(doc_map_path, 'r') as f:
                self.document_map = json.load(f)
            
            st.sidebar.success(f"📚 Loaded vector store: {self.index.ntotal} vectors from {len(self.document_map)} documents")
            
        except Exception as e:
            st.sidebar.warning(f"⚠️ Error loading vector store: {str(e)}")
            st.sidebar.info("🔄 Starting with fresh vector store")
            self.index = faiss.IndexFlatIP(self.dimension)
            self.metadata = []
            self.document_map = {}

# LLM Handler
class BaseLLM(ABC):
    """Abstract base class for LLM implementations."""
    
    @abstractmethod
    def generate_response(self, query: str, context: str, **kwargs) -> str:
        pass
    
    @abstractmethod
    def get_model_info(self) -> Dict[str, Any]:
        pass

class GroqLLM(BaseLLM):
    """Professional Groq API implementation with streaming support."""
    
    SUPPORTED_MODELS = {
        "openai/gpt-oss-120b": "OpenAI GPT-OSS 120B",
        "llama-3.1-70b-versatile": "LLaMA 3.1 70B Versatile",
        "llama-3.1-8b-instant": "LLaMA 3.1 8B Instant",
        "mixtral-8x7b-32768": "Mixtral 8x7B",
        "gemma2-9b-it": "Gemma2 9B IT"
    }

    def __init__(self, api_key: Optional[str] = None, model: str = "openai/gpt-oss-120b"):
        self.api_key = api_key or os.getenv("GROQ_API_KEY")
        self.model = model

        if not self.api_key:
            raise ValueError("Groq API key required. Set GROQ_API_KEY environment variable.")

        if model not in self.SUPPORTED_MODELS:
            st.warning(f"⚠️ Unknown model {model}, using default")
            self.model = "openai/gpt-oss-120b"

        self.client = Groq(api_key=self.api_key)
        st.sidebar.success(f"✅ Initialized Groq LLM with {self.SUPPORTED_MODELS.get(self.model, self.model)}")
    
    def generate_response(self, query: str, context: str, 
                         max_tokens: int = 8192,
                         temperature: float = 1.0, 
                         top_p: float = 1.0,
                         reasoning_effort: str = "medium", 
                         enable_streaming: bool = True) -> str:
        """Generate response using Groq API with advanced parameters."""

        system_prompt = """You are an expert AI assistant specializing in document analysis and knowledge extraction. 

Your responsibilities:
- Provide accurate, well-structured answers based solely on the provided context
- Cite specific information from the context when possible
- Clearly state when information is insufficient to answer the question
- Maintain professional, concise communication
- Focus on factual accuracy over speculation"""

        user_prompt = f"""Context Information:
{context}

User Question: {query}

Please provide a comprehensive answer based on the context above. If the context doesn't contain sufficient information, clearly state this limitation."""

        try:
            completion = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=temperature,
                max_completion_tokens=max_tokens,
                top_p=top_p,
                reasoning_effort=reasoning_effort,
                stream=enable_streaming,
                stop=None
            )

            if enable_streaming:
                response_text = ""
                for chunk in completion:
                    if chunk.choices[0].delta.content:
                        response_text += chunk.choices[0].delta.content
                return response_text.strip()
            else:
                return completion.choices[0].message.content.strip()

        except Exception as e:
            return f"❌ Error generating response: {str(e)}"
    
    def get_model_info(self) -> Dict[str, Any]:
        """Get Groq model information."""
        return {
            'provider': 'groq',
            'model': self.model,
            'model_name': self.SUPPORTED_MODELS.get(self.model, self.model),
            'supports_streaming': True,
            'max_tokens': 8192
        }

class LLMHandler:
    """Professional LLM handler with intelligent provider management."""

    def __init__(self, llm_type: str = "groq", **kwargs):
        self.llm_type = llm_type

        if llm_type == "groq":
            self.llm = GroqLLM(**kwargs)
        else:
            raise ValueError(f"Unsupported LLM type: {llm_type}")
    
    def generate_answer(self,
                       query: str,
                       retrieved_chunks: List[Dict[str, Any]],
                       max_tokens: int = 8192,
                       **kwargs) -> Dict[str, Any]:
        if not retrieved_chunks:
            return {
                'answer': "I couldn't find any relevant information in the knowledge base to answer your question.",
                'sources': [],
                'context_used': "",
                'num_sources': 0,
                'llm_type': self.llm_type,
                'model_used': getattr(self.llm, 'model', 'unknown')
            }
        
        # Prepare context from top chunks
        context_parts = []
        sources = []
        
        for i, chunk in enumerate(retrieved_chunks[:5]):
            context_parts.append(f"[Source {i+1}] {chunk['text']}")
            sources.append({
                'file_name': chunk['file_name'],
                'file_path': chunk['file_path'],
                'similarity_score': chunk.get('similarity_score', 0.0),
                'chunk_id': chunk.get('chunk_id', 0)
            })
        
        context = "\n\n".join(context_parts)
        
        # Generate answer
        answer = self.llm.generate_response(
            query, context, max_tokens=max_tokens, **kwargs
        )
        
        return {
            'answer': answer,
            'sources': sources,
            'context_used': context,
            'num_sources': len(sources),
            'llm_type': self.llm_type,
            'model_used': getattr(self.llm, 'model', 'unknown')
        }
    
    def get_model_info(self) -> Dict[str, Any]:
        """Get comprehensive model information."""
        return self.llm.get_model_info()

def create_llm_handler(llm_type: str = "auto", **kwargs) -> LLMHandler:
    """Create LLM handler with intelligent provider selection."""
    if llm_type == "auto":
        # Try Groq first
        if os.getenv("GROQ_API_KEY"):
            try:
                return LLMHandler("groq", **kwargs)
            except Exception as e:
                st.warning(f"⚠️ Failed to initialize Groq LLM: {str(e)}")

    return LLMHandler("groq", **kwargs)

# ==================== STREAMLIT APP ====================

# Streamlit configuration
st.set_page_config(
    page_title="LocalRAG AI Knowledge Assistant",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://github.com/localrag/localrag-ai',
        'Report a bug': 'https://github.com/localrag/localrag-ai/issues',
        'About': """
        # 🔍 LocalRAG AI Knowledge Assistant
        
        **Academic Research Project**
        
        A Retrieval-Augmented Generation (RAG) system for intelligent document analysis 
        and question answering using state-of-the-art AI models.
        """
    }
)

# Initialize session state
if 'initialized' not in st.session_state:
    st.session_state.initialized = False
    st.session_state.vector_store = None
    st.session_state.llm_handler = None
    st.session_state.document_processor = None

@st.cache_resource
def initialize_system():
    """Initialize the LocalRAG system components."""
    try:
        # Initialize document processor
        document_processor = DocumentProcessor(chunk_size=1000, chunk_overlap=200)
        
        # Initialize vector store
        store_path = os.getenv("VECTOR_STORE_PATH", "data/vector_store")
        embedding_model = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
        vector_store = FAISSVectorStore(store_path=store_path, embedding_model=embedding_model)
        
        # Initialize LLM handler
        llm_handler = create_llm_handler(llm_type="auto")
        
        return document_processor, vector_store, llm_handler
    except Exception as e:
        st.error(f"Failed to initialize system: {str(e)}")
        return None, None, None

def main():
    """Main Streamlit application."""
    
    # Header
    st.markdown("""
    <div style="text-align: center; padding: 20px 0;">
        <h1 style="color: #1f77b4; margin-bottom: 0; font-size: 2.5rem;">
            🔍 LocalRAG AI Knowledge Assistant
        </h1>
        <p style="font-size: 1.2rem; color: #666; margin-top: 5px;">
            Academic Research Project - Standalone Version
        </p>
        <hr style="margin: 20px 0; border: 1px solid #e0e0e0;">
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style="text-align: center; margin-bottom: 30px;">
        <p style="font-size: 1rem; color: #555;">
            📁 Upload Research Documents • 🤖 Ask Intelligent Questions • 📊 Get AI-Powered Insights
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Initialize system
    if not st.session_state.initialized:
        with st.spinner("🚀 Initializing LocalRAG AI Knowledge Assistant..."):
            document_processor, vector_store, llm_handler = initialize_system()
            
            if all([document_processor, vector_store, llm_handler]):
                st.session_state.document_processor = document_processor
                st.session_state.vector_store = vector_store
                st.session_state.llm_handler = llm_handler
                st.session_state.initialized = True
                st.success("✅ System initialized successfully!")
            else:
                st.error("❌ Failed to initialize system. Please check your configuration.")
                return
    
    # Sidebar
    with st.sidebar:
        st.header("📊 System Dashboard")
        
        # System stats
        if st.session_state.initialized:
            stats = st.session_state.vector_store.get_stats()
            
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Documents", stats.get("total_documents", 0))
                st.metric("Vectors", stats.get("total_vectors", 0))
            with col2:
                st.metric("Storage (MB)", stats.get("storage_size_mb", 0))
                model_info = st.session_state.llm_handler.get_model_info()
                model_name = model_info.get('model', 'Unknown')[:15] + "..."
                st.metric("Model", model_name)
            
            st.info(f"🤖 **Provider:** {model_info.get('provider', 'Unknown')}")
            st.info(f"🔤 **Embeddings:** {stats.get('embedding_model', 'Unknown')}")
        else:
            st.warning("System not initialized")
        
        st.divider()
        
        # System Management
        st.subheader("🛠️ System Management")
        
        if st.button("🔄 Refresh Stats", type="secondary"):
            st.rerun()
        
        if st.button("🗑️ Clear Knowledge Base", type="secondary"):
            if st.session_state.get("confirm_clear", False):
                if st.session_state.initialized:
                    st.session_state.vector_store.clear_store()
                    st.success("✅ Knowledge base cleared!")
                    st.rerun()
                st.session_state.confirm_clear = False
            else:
                st.session_state.confirm_clear = True
                st.warning("⚠️ Click again to confirm")
    
    # Main tabs
    tab1, tab2, tab3 = st.tabs([
        "📁 Document Ingestion",
        "🔍 Intelligent Query",
        "📊 Knowledge Analytics"
    ])
    
    with tab1:
        st.header("📁 Document Ingestion Center")
        st.markdown("*Process and index documents for intelligent analysis*")
        
        if not st.session_state.initialized:
            st.warning("Please wait for system initialization to complete.")
            return
        
        # File upload
        uploaded_files = st.file_uploader(
            "Select files to upload and process",
            type=['pdf', 'txt', 'docx', 'csv', 'xlsx', 'md', 'pptx'],
            accept_multiple_files=True,
            help="Supported: PDF, TXT, DOCX, CSV, XLSX, MD, PPTX"
        )
        
        if uploaded_files:
            st.write(f"**Selected {len(uploaded_files)} files:**")
            for file in uploaded_files:
                st.write(f"- {file.name} ({file.size:,} bytes)")
            
            if st.button("🚀 Process Files", type="primary"):
                processed_files = []
                failed_files = []
                total_chunks = 0
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_path = Path(temp_dir)
                    
                    for i, file in enumerate(uploaded_files):
                        try:
                            status_text.text(f"Processing {file.name}...")
                            progress_bar.progress((i + 1) / len(uploaded_files))
                            
                            # Save uploaded file
                            file_path = temp_path / file.name
                            with open(file_path, "wb") as f:
                                f.write(file.getvalue())
                            
                            # Process the file
                            doc_data = st.session_state.document_processor.process_file(file_path)
                            
                            # Add to vector store
                            st.session_state.vector_store.add_documents([doc_data])
                            
                            processed_files.append(file.name)
                            total_chunks += doc_data['num_chunks']
                            
                        except Exception as e:
                            failed_files.append(f"{file.name}: {str(e)}")
                
                status_text.empty()
                progress_bar.empty()
                
                if processed_files:
                    st.success(f"✅ Successfully processed {len(processed_files)} files")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.metric("Files Processed", len(processed_files))
                    with col2:
                        st.metric("Chunks Created", total_chunks)
                    
                    with st.expander("✅ Successfully Processed"):
                        for file in processed_files:
                            st.write(f"- {file}")
                
                if failed_files:
                    st.error("❌ Some files failed to process")
                    with st.expander("❌ Failed to Process"):
                        for file in failed_files:
                            st.write(f"- {file}")
    
    with tab2:
        st.header("🔍 Intelligent Query Interface")
        st.markdown("*Ask questions and get AI-powered answers from your knowledge base*")
        
        if not st.session_state.initialized:
            st.warning("Please wait for system initialization to complete.")
            return
        
        # Query input
        query = st.text_area(
            "**Your Question:**",
            placeholder="Ask anything about your documents...",
            height=100,
            help="Enter your question in natural language"
        )
        
        # Advanced settings
        with st.expander("⚙️ Advanced Query Settings"):
            col1, col2 = st.columns(2)
            
            with col1:
                k = st.slider("Number of Sources", 1, 10, 5, help="How many relevant documents to retrieve")
                score_threshold = st.slider("Similarity Threshold", 0.0, 1.0, 0.0, 0.1, help="Minimum similarity score")
            
            with col2:
                temperature = st.slider("Response Creativity", 0.0, 2.0, 1.0, 0.1, help="Higher = more creative")
                max_tokens = st.slider("Max Response Length", 100, 8192, 2000, 100, help="Maximum response tokens")
        
        if query and st.button("🔍 Search & Analyze", type="primary"):
            with st.spinner("🤖 AI is analyzing your question..."):
                try:
                    # Search vector store
                    search_results = st.session_state.vector_store.search(
                        query=query,
                        k=k,
                        score_threshold=score_threshold
                    )
                    
                    # Generate AI answer
                    result = st.session_state.llm_handler.generate_answer(
                        query=query,
                        retrieved_chunks=search_results,
                        max_tokens=max_tokens,
                        temperature=temperature,
                        top_p=1.0,
                        reasoning_effort="medium",
                        enable_streaming=True
                    )
                    
                    # Display answer
                    st.subheader("🤖 AI Response")
                    st.markdown(result['answer'])
                    
                    # Display sources
                    if result['sources']:
                        st.subheader(f"📚 Sources ({result['num_sources']})")
                        
                        for i, source in enumerate(result['sources'], 1):
                            with st.expander(f"Source {i}: {source['file_name']} (Score: {source['similarity_score']:.3f})"):
                                st.write(f"**File:** {source['file_name']}")
                                st.write(f"**Path:** {source['file_path']}")
                                st.write(f"**Similarity:** {source['similarity_score']:.3f}")
                                st.write(f"**Chunk ID:** {source['chunk_id']}")
                    
                    # Display metadata
                    with st.expander("🔍 Query Metadata"):
                        col1, col2 = st.columns(2)
                        with col1:
                            st.write(f"**LLM Provider:** {result['llm_type']}")
                            st.write(f"**Model Used:** {result['model_used']}")
                        with col2:
                            st.write(f"**Sources Found:** {result['num_sources']}")
                            st.write(f"**Query:** {query}")
                            
                except Exception as e:
                    st.error(f"❌ Error processing query: {str(e)}")
    
    with tab3:
        st.header("📊 Knowledge Base Analytics")
        st.markdown("*Comprehensive insights into your document collection*")
        
        if not st.session_state.initialized:
            st.warning("Please wait for system initialization to complete.")
            return
        
        stats = st.session_state.vector_store.get_stats()
        
        # Overview metrics
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric(
                "Total Documents", 
                stats.get("total_documents", 0),
                help="Number of documents in knowledge base"
            )
        
        with col2:
            st.metric(
                "Vector Embeddings", 
                stats.get("total_vectors", 0),
                help="Total text chunks indexed"
            )
        
        with col3:
            st.metric(
                "Storage Size", 
                f"{stats.get('storage_size_mb', 0)} MB",
                help="Disk space used by vector store"
            )
        
        with col4:
            avg_chunks = (stats.get("total_vectors", 0) / max(stats.get("total_documents", 1), 1))
            st.metric(
                "Avg Chunks/Doc", 
                f"{avg_chunks:.1f}",
                help="Average chunks per document"
            )
        
        st.divider()
        
        # Document list
        if stats.get("documents"):
            st.subheader("📁 Document Collection")
            
            for i, doc_path in enumerate(stats["documents"], 1):
                doc_name = Path(doc_path).name
                st.write(f"{i}. **{doc_name}**")
                st.caption(f"Path: {doc_path}")
        else:
            st.info("📭 No documents in knowledge base yet. Upload some documents to get started!")
        
        # System information
        st.divider()
        st.subheader("🛠️ System Configuration")
        
        model_info = st.session_state.llm_handler.get_model_info()
        
        col1, col2 = st.columns(2)
        with col1:
            st.write(f"**Embedding Model:** {stats.get('embedding_model', 'Unknown')}")
            st.write(f"**LLM Provider:** {model_info.get('provider', 'Unknown')}")
        with col2:
            st.write(f"**Current Model:** {model_info.get('model', 'Unknown')}")
            st.write(f"**System Version:** 2.0.0")

if __name__ == "__main__":
    main()
